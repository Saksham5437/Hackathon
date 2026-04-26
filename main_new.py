"""
╔══════════════════════════════════════════════════════════════════╗
║          MINI NOTEBOOKLM — Backend (main.py) — FIXED             ║
║         FastAPI + RAG + ChromaDB + Gemini/OpenAI API             ║
╚══════════════════════════════════════════════════════════════════╝
"""

import os
import uuid
import shutil
import logging
import re
import requests
import json
from pathlib import Path
from typing import Optional, List
from datetime import datetime
from collections import deque

import uvicorn
from dotenv import load_dotenv
from fastapi import FastAPI, File, UploadFile, HTTPException, Query, Header
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

# Document loaders
from pypdf import PdfReader
import docx  # python-docx
from pptx import Presentation  # python-pptx

# Text-to-speech
from gtts import gTTS

# LangChain components
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_core.documents import Document

# ─────────────────────────────────────────────
# CONFIG
# ─────────────────────────────────────────────

load_dotenv()

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
logger = logging.getLogger("mini-notebooklm")

UPLOAD_DIR = Path("uploads")
CHROMA_DIR = Path("chroma_db")
VOICE_DIR  = Path("voice_overviews")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
CHROMA_DIR.mkdir(parents=True, exist_ok=True)
VOICE_DIR.mkdir(parents=True, exist_ok=True)
USERS_FILE = Path("users.json")

from passlib.context import CryptContext
pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")

def load_users() -> dict:
    if not USERS_FILE.exists():
        return {}
    try:
        return json.loads(USERS_FILE.read_text())
    except:
        return {}

def save_users(users: dict):
    USERS_FILE.write_text(json.dumps(users, indent=2))

def validate_password(password: str) -> Optional[str]:
    if len(password) > 15:
        return "Password must be maximum 15 characters."
    if not any(c.isupper() for c in password):
        return "Password must contain at least one uppercase letter."
    if not any(c.islower() for c in password):
        return "Password must contain at least one lowercase letter."
    if not any(c.isdigit() for c in password):
        return "Password must contain at least one number."
    if not any(c in "!@#$%^&*()_+-=[]{}|;:,.<>?" for c in password):
        return "Password must contain at least one special character."
    return None

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
LLM_PROVIDER   = os.getenv("LLM_PROVIDER", "gemini").lower()
YOUTUBE_API_KEY = os.getenv("YOUTUBE_API_KEY", "")
EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL", "all-MiniLM-L6-v2")
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")

CHUNK_SIZE    = int(os.getenv("CHUNK_SIZE", "800"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "100"))
TOP_K_CHUNKS  = int(os.getenv("TOP_K_CHUNKS", "5"))
MAX_HISTORY   = int(os.getenv("MAX_HISTORY", "10"))
LLM_MAX_OUTPUT_TOKENS = int(os.getenv("LLM_MAX_OUTPUT_TOKENS", "3072"))

HOST = os.getenv("HOST", "0.0.0.0")
PORT = int(os.getenv("PORT", "5000"))

ALLOWED_EXTENSIONS = {".pdf", ".txt", ".docx", ".pptx"}

# ─────────────────────────────────────────────
# APP INIT
# ─────────────────────────────────────────────

app = FastAPI(
    title="Mini NotebookLM API",
    description="A lightweight RAG-powered document chat backend.",
    version="1.1.0",
    docs_url="/docs",
    redoc_url="/redoc",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ─────────────────────────────────────────────
# PYDANTIC SCHEMAS
# ─────────────────────────────────────────────

class AskRequest(BaseModel):
    question: str
    session_id: Optional[str] = "default"
    include_sources: Optional[bool] = True
    file_name: Optional[str] = None

class AskResponse(BaseModel):
    answer: str
    sources: Optional[List[dict]] = []
    session_id: str

class SummarizeRequest(BaseModel):
    file_name: Optional[str] = None

class SummarizeResponse(BaseModel):
    summary: str
    file_name: Optional[str] = None

class DeleteResponse(BaseModel):
    message: str
    file_name: str

class HealthResponse(BaseModel):
    status: str
    provider: str
    uploaded_files: int
    timestamp: str

class VoiceOverviewRequest(BaseModel):
    file_name: Optional[str] = None
    language: Optional[str] = "en"

class VoiceOverviewResponse(BaseModel):
    audio_file: str
    summary_text: str
    file_name: Optional[str] = None

class ConceptMapRequest(BaseModel):
    file_name: Optional[str] = None
    output_format: Optional[str] = "mermaid"

class ConceptMapResponse(BaseModel):
    concept_map: str
    output_format: str
    file_name: Optional[str] = None

class AuthRequest(BaseModel):
    username: str
    password: str

class YouTubeSearchRequest(BaseModel):
    topic: Optional[str] = None
    file_name: Optional[str] = None
    language_code: Optional[str] = None
    max_duration_minutes: Optional[int] = None
    sort_by: Optional[str] = "views"
    max_results: Optional[int] = 5

class YouTubeVideoItem(BaseModel):
    title: str
    channel: str
    url: str
    thumbnail: str
    duration_seconds: int
    duration_label: str
    view_count: int
    like_count: int
    published_at: str
    description_snippet: str

class YouTubeSearchResponse(BaseModel):
    topic: str
    videos: List[YouTubeVideoItem]
    total_returned: int

# ─────────────────────────────────────────────
# SESSION MEMORY
# ─────────────────────────────────────────────

session_store: dict[str, deque] = {}

def get_session_history(session_id: str) -> list[dict]:
    if session_id not in session_store:
        session_store[session_id] = deque(maxlen=MAX_HISTORY)
    return list(session_store[session_id])

def append_to_session(session_id: str, role: str, content: str):
    if session_id not in session_store:
        session_store[session_id] = deque(maxlen=MAX_HISTORY)
    session_store[session_id].append({"role": role, "content": content})

# ─────────────────────────────────────────────
# EMBEDDINGS & VECTOR STORE
# ─────────────────────────────────────────────

logger.info(f"Loading embedding model: {EMBEDDING_MODEL}")
embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)

vector_store = Chroma(
    persist_directory=str(CHROMA_DIR),
    embedding_function=embeddings,
    collection_name="mini_notebooklm",
)

def get_retriever(k: int = TOP_K_CHUNKS, file_name: Optional[str] = None, user_id: Optional[str] = None):
    search_kwargs = {"k": k}
    filter_dict = {}
    if user_id:
        filter_dict["user"] = user_id
    if file_name:
        filter_dict["source"] = file_name
    
    if filter_dict:
        if len(filter_dict) > 1:
            search_kwargs["filter"] = {"$and": [{k: v} for k, v in filter_dict.items()]}
        else:
            search_kwargs["filter"] = filter_dict
            
    return vector_store.as_retriever(search_kwargs=search_kwargs)

# ─────────────────────────────────────────────
# LLM CALLER - Gemini + OpenAI
# ─────────────────────────────────────────────

def call_llm(prompt: str) -> str:
    """Send prompt to LLM and return response. Raises HTTPException on failure."""
    if LLM_PROVIDER == "openai":
        if not OPENAI_API_KEY:
            raise HTTPException(status_code=500, detail="OPENAI_API_KEY not set in .env")
        try:
            from openai import OpenAI
            client = OpenAI(api_key=OPENAI_API_KEY)
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.2,
                max_tokens=LLM_MAX_OUTPUT_TOKENS,
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            logger.error(f"OpenAI API error: {e}")
            raise HTTPException(status_code=500, detail=f"OpenAI API error: {str(e)}")

    elif LLM_PROVIDER == "gemini":
        if not GEMINI_API_KEY:
            raise HTTPException(status_code=500, detail="GEMINI_API_KEY not set in .env")
        try:
            response = requests.post(
                f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent",
                headers={
                    "x-goog-api-key": GEMINI_API_KEY,
                    "Content-Type": "application/json",
                },
                json={
                    "contents": [
                        {
                            "role": "user",
                            "parts": [{"text": prompt}],
                        }
                    ],
                    "generationConfig": {
                        "temperature": 0.2,
                        "maxOutputTokens": LLM_MAX_OUTPUT_TOKENS,
                    },
                },
                timeout=45,
            )
            response.raise_for_status()
            data = response.json()
            parts = data.get("candidates", [{}])[0].get("content", {}).get("parts", [])
            text = "".join(part.get("text", "") for part in parts).strip()
            if not text:
                raise ValueError("Gemini returned an empty response.")
            return text
        except requests.HTTPError as e:
            detail = e.response.text if e.response is not None else str(e)
            logger.error(f"Gemini API error: {detail}")
            raise HTTPException(status_code=500, detail=f"Gemini API error: {detail}")
        except Exception as e:
            logger.error(f"Gemini API error: {e}")
            raise HTTPException(status_code=500, detail=f"Gemini API error: {str(e)}")

    else:
        raise HTTPException(status_code=500, detail=f"Unknown LLM_PROVIDER: '{LLM_PROVIDER}'. Set to 'gemini' or 'openai' in .env")

# ─────────────────────────────────────────────
# TEXT EXTRACTION
# ─────────────────────────────────────────────

def extract_text_from_pdf(file_path: Path) -> str:
    reader = PdfReader(str(file_path))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages)

def extract_text_from_txt(file_path: Path) -> str:
    return file_path.read_text(encoding="utf-8", errors="ignore")

def extract_text_from_docx(file_path: Path) -> str:
    doc = docx.Document(str(file_path))
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())

def extract_text_from_pptx(file_path: Path) -> str:
    prs = Presentation(str(file_path))
    lines = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)

def extract_text(file_path: Path) -> str:
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".txt":
        return extract_text_from_txt(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

# ─────────────────────────────────────────────
# CHUNKING & INDEXING
# ─────────────────────────────────────────────

def chunk_and_index(text: str, file_name: str, user_id: str) -> int:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", ".", " ", ""],
    )
    chunks = splitter.split_text(text)
    documents = [
        Document(
            page_content=chunk,
            metadata={"source": file_name, "chunk_index": i, "user": user_id},
        )
        for i, chunk in enumerate(chunks)
    ]
    vector_store.add_documents(documents)
    logger.info(f"Indexed {len(documents)} chunks from '{file_name}' for user '{user_id}'")
    return len(documents)


@app.post("/register", summary="Register a new user")
async def register_user(payload: AuthRequest):
    users = load_users()
    username = payload.username.strip()
    if username in users:
        raise HTTPException(status_code=400, detail="Username already exists. Please choose another or sign in.")
    
    error = validate_password(payload.password)
    if error:
        raise HTTPException(status_code=400, detail=error)
    
    hashed_pwd = pwd_context.hash(payload.password)
    users[username] = hashed_pwd
    save_users(users)
    
    return {"message": "User registered successfully.", "username": username}


@app.post("/login", summary="Login an existing user")
async def login_user(payload: AuthRequest):
    users = load_users()
    username = payload.username.strip()
    if username not in users:
        raise HTTPException(status_code=401, detail="Username not found. Please register first.")
    
    if not pwd_context.verify(payload.password, users[username]):
        raise HTTPException(status_code=401, detail="Incorrect password. Please try again.")
    
    return {"message": "Login successful.", "username": username}

# ─────────────────────────────────────────────
# RAG PROMPT BUILDER
# ─────────────────────────────────────────────

def build_rag_prompt(question: str, context_chunks: list[str], history: list[dict]) -> str:
    context = "\n\n---\n\n".join(context_chunks)

    history_text = ""
    if history:
        history_lines = []
        for msg in history[-6:]:
            role = "User" if msg["role"] == "user" else "Assistant"
            history_lines.append(f"{role}: {msg['content']}")
        history_text = "\n".join(history_lines)

    prompt = f"""You are Aura, a NotebookLM-style research assistant. Answer the user's question using ONLY the document context below.

Core rules:
- Give a complete, well-organized answer. Do not stop early.
- Use clear Markdown with short headings and bullets when they improve readability.
- Start with a direct answer, then explain the supporting details.
- Ground claims in the provided context. Mention source names naturally when useful.
- If the context is incomplete, say what is missing instead of guessing.
- If the answer is not in the context, say exactly: "I couldn't find relevant information in the uploaded documents."
- Prefer this structure when it fits the question: Short answer, Key points, Explanation, Evidence from the document, and Limitations or what to ask next.
- Do not include irrelevant boilerplate or generic textbook material unless the document supports it.

━━━━━━━━━━━━━ DOCUMENT CONTEXT ━━━━━━━━━━━━━
{context}

━━━━━━━━━━━━━ RECENT CONVERSATION ━━━━━━━━━━
{history_text if history_text else "(No prior conversation)"}

━━━━━━━━━━━━━ CURRENT QUESTION ━━━━━━━━━━━━
{question}

Answer:"""
    return prompt

# ─────────────────────────────────────────────
# ROUTES
# ─────────────────────────────────────────────

@app.post("/upload", summary="Upload a document (PDF, TXT, DOCX, PPTX)")
async def upload_file(
    file: UploadFile = File(...),
    x_user_profile: str = Header(..., alias="X-User-Profile")
):
    user_id = x_user_profile.strip()
    user_dir = UPLOAD_DIR / user_id
    user_dir.mkdir(parents=True, exist_ok=True)

    file_ext = Path(file.filename).suffix.lower()
    if file_ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type '{file_ext}'. Allowed: {', '.join(ALLOWED_EXTENSIONS)}",
        )

    save_path = user_dir / file.filename
    try:
        with open(save_path, "wb") as f:
            shutil.copyfileobj(file.file, f)
        logger.info(f"Saved: {save_path} for user {user_id}")
    except Exception as e:
        logger.error(f"Failed to save file: {e}")
        raise HTTPException(status_code=500, detail=f"File save failed: {str(e)}")

    try:
        text = extract_text(save_path)
        if not text.strip():
            raise ValueError("Extracted text is empty. The file may be blank or image-only.")
    except Exception as e:
        logger.error(f"Text extraction failed: {e}")
        save_path.unlink(missing_ok=True)
        raise HTTPException(status_code=422, detail=f"Text extraction failed: {str(e)}")

    try:
        num_chunks = chunk_and_index(text, file.filename, user_id)
    except Exception as e:
        logger.error(f"Indexing failed: {e}")
        raise HTTPException(status_code=500, detail=f"Indexing failed: {str(e)}")

    return JSONResponse(
        status_code=200,
        content={
            "message": f"File '{file.filename}' uploaded and indexed successfully.",
            "file_name": file.filename,
            "chunks_indexed": num_chunks,
            "file_size_bytes": save_path.stat().st_size,
        },
    )


@app.post("/ask", response_model=AskResponse, summary="Ask a question about your documents")
async def ask_question(
    payload: AskRequest,
    x_user_profile: str = Header(..., alias="X-User-Profile")
):
    user_id = x_user_profile.strip()
    user_dir = UPLOAD_DIR / user_id

    """FIXED: Properly handles all error cases and returns clean answer."""
    question = payload.question.strip()
    if not question:
        raise HTTPException(status_code=400, detail="Question cannot be empty.")

    session_id = payload.session_id or "default"
    file_name = payload.file_name

    # Check if any documents exist first
    try:
        if file_name and not (user_dir / file_name).exists():
            raise HTTPException(status_code=404, detail=f"File '{file_name}' not found for your profile.")
        all_docs = vector_store.get(where={"user": user_id})
        if not all_docs or not all_docs.get("documents"):
            return AskResponse(
                answer="No documents have been uploaded yet. Please upload a PDF, TXT, DOCX, or PPTX file using the sidebar first.",
                sources=[],
                session_id=session_id,
            )
    except Exception as e:
        logger.warning(f"Could not check vector store: {e}")

    # Retrieve relevant chunks
    try:
        retriever = get_retriever(k=TOP_K_CHUNKS, file_name=file_name, user_id=user_id)
        relevant_docs = retriever.invoke(question)
    except Exception as e:
        logger.error(f"Retrieval error: {e}")
        raise HTTPException(status_code=500, detail=f"Document retrieval failed: {str(e)}")

    if not relevant_docs:
        return AskResponse(
            answer="I couldn't find relevant information in the uploaded documents for your question. Try rephrasing or uploading more relevant documents.",
            sources=[],
            session_id=session_id,
        )

    context_chunks = [doc.page_content for doc in relevant_docs]
    sources = [
        {
            "source": doc.metadata.get("source", "unknown"),
            "chunk_index": doc.metadata.get("chunk_index", -1),
            "preview": doc.page_content[:200] + "..." if len(doc.page_content) > 200 else doc.page_content,
        }
        for doc in relevant_docs
    ]

    history = get_session_history(session_id)
    prompt = build_rag_prompt(question, context_chunks, history)

    # FIXED: call_llm now raises proper HTTPException on failure
    answer = call_llm(prompt)

    # Persist to session memory
    append_to_session(session_id, "user", question)
    append_to_session(session_id, "assistant", answer)

    logger.info(f"[{session_id}] Q: {question[:60]}... → A: {answer[:60]}...")

    return AskResponse(
        answer=answer,
        sources=sources if payload.include_sources else [],
        session_id=session_id,
    )


@app.post("/summarize", response_model=SummarizeResponse, summary="Summarize uploaded documents")
async def summarize(
    payload: SummarizeRequest,
    x_user_profile: str = Header(..., alias="X-User-Profile")
):
    user_id = x_user_profile.strip()
    user_dir = UPLOAD_DIR / user_id
    if payload.file_name:
        file_path = user_dir / payload.file_name
        if not file_path.exists():
            raise HTTPException(status_code=404, detail=f"File '{payload.file_name}' not found for your profile.")
        try:
            text = extract_text(file_path)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Text extraction failed: {str(e)}")
        text_sample = text[:12000]
        label = payload.file_name
    else:
        all_docs = vector_store.get(where={"user": user_id})
        if not all_docs or not all_docs.get("documents"):
            raise HTTPException(
                status_code=404,
                detail="No documents indexed yet for your profile. Please upload a file first.",
            )
        text_sample = "\n\n".join(all_docs["documents"][:30])[:12000]
        label = "all uploaded documents for your profile"

    prompt = f"""You are a document summarization expert.
Provide a clear, structured, and concise summary of the following document content.
Include: main topics, key points, and any notable conclusions.
Keep the summary under 350 words. Use plain paragraphs, no bullet points or markdown symbols.

DOCUMENT CONTENT:
{text_sample}

SUMMARY:"""

    summary = call_llm(prompt)
    logger.info(f"Summarized: {label}")
    return SummarizeResponse(summary=summary, file_name=payload.file_name)


@app.get("/files", summary="List all uploaded files")
async def list_files(x_user_profile: str = Header(..., alias="X-User-Profile")):
    user_id = x_user_profile.strip()
    user_dir = UPLOAD_DIR / user_id
    user_dir.mkdir(parents=True, exist_ok=True)
    
    files = []
    for f in sorted(user_dir.iterdir()):
        if f.is_file():
            stat = f.stat()
            files.append({
                "file_name": f.name,
                "size_bytes": stat.st_size,
                "uploaded_at": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "extension": f.suffix.lower(),
            })
    return {"files": files, "total": len(files)}


@app.delete("/files/{file_name}", response_model=DeleteResponse, summary="Delete an uploaded file")
async def delete_file(
    file_name: str,
    x_user_profile: str = Header(..., alias="X-User-Profile")
):
    user_id = x_user_profile.strip()
    user_dir = UPLOAD_DIR / user_id
    file_path = user_dir / file_name
    if not file_path.exists():
        raise HTTPException(status_code=404, detail=f"File '{file_name}' not found for your profile.")
    try:
        file_path.unlink()
        logger.info(f"Deleted file: {file_name}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to delete: {str(e)}")
    return DeleteResponse(message=f"File '{file_name}' deleted.", file_name=file_name)


@app.get("/session/{session_id}", summary="Get chat history for a session")
async def get_session(session_id: str):
    history = get_session_history(session_id)
    return {"session_id": session_id, "messages": history, "count": len(history)}


@app.delete("/session/{session_id}", summary="Clear chat history for a session")
async def clear_session(session_id: str):
    if session_id in session_store:
        del session_store[session_id]
    return {"message": f"Session '{session_id}' cleared.", "session_id": session_id}


@app.get("/health", response_model=HealthResponse, summary="Server health check")
async def health_check(x_user_profile: str = Header(..., alias="X-User-Profile")):
    user_id = x_user_profile.strip()
    user_dir = UPLOAD_DIR / user_id
    user_dir.mkdir(parents=True, exist_ok=True)
    uploaded_count = sum(1 for f in user_dir.iterdir() if f.is_file())
    return HealthResponse(
        status="ok",
        provider=LLM_PROVIDER,
        uploaded_files=uploaded_count,
        timestamp=datetime.utcnow().isoformat() + "Z",
    )


@app.get("/", include_in_schema=False)
async def root():
    """Serve the frontend HTML interface."""
    frontend_path = Path(__file__).parent / "index.html"
    if frontend_path.exists():
        return FileResponse(str(frontend_path), media_type="text/html")
    return JSONResponse(content={
        "message": "Mini NotebookLM API is running! Place index.html in the same folder to serve the UI.",
        "docs": f"http://localhost:{PORT}/docs",
        "health": f"http://localhost:{PORT}/health",
    })


# ─────────────────────────────────────────────
# VOICE OVERVIEW
# ─────────────────────────────────────────────

@app.post("/voice-overview", summary="Generate a spoken audio overview of your documents")
async def voice_overview(
    payload: VoiceOverviewRequest,
    x_user_profile: str = Header(..., alias="X-User-Profile")
):
    user_id = x_user_profile.strip()
    user_dir = UPLOAD_DIR / user_id
    if payload.file_name:
        file_path = user_dir / payload.file_name
        if not file_path.exists():
            raise HTTPException(status_code=404, detail=f"File '{payload.file_name}' not found for your profile.")
        try:
            text = extract_text(file_path)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Text extraction failed: {str(e)}")
        text_sample = text[:12000]
        label = payload.file_name
    else:
        all_docs = vector_store.get(where={"user": user_id})
        if not all_docs or not all_docs.get("documents"):
            raise HTTPException(
                status_code=404,
                detail="No documents indexed yet for your profile. Please upload a file first.",
            )
        text_sample = "\n\n".join(all_docs["documents"][:30])[:12000]
        label = "all uploaded documents for your profile"

    prompt = f"""You are creating a friendly spoken-audio overview for a student or professional.
Write a clear, engaging, conversational summary of the following document content.
The summary should sound completely natural when read aloud — no bullet points, no markdown, no special characters.
Speak in flowing sentences and paragraphs only. Keep it under 200 words so it is easy to listen to.

DOCUMENT CONTENT:
{text_sample}

SPOKEN SUMMARY:"""

    summary_text = call_llm(prompt)

    audio_filename = f"overview_{uuid.uuid4().hex[:8]}.mp3"
    audio_path = VOICE_DIR / audio_filename

    try:
        # Clean text for TTS — remove any markdown remnants
        clean_text = re.sub(r'[*_#`]', '', summary_text)
        clean_text = re.sub(r'\n+', ' ', clean_text).strip()

        tts = gTTS(text=clean_text, lang=payload.language or "en", slow=False)
        tts.save(str(audio_path))
        logger.info(f"Saved voice overview: {audio_path}")
    except Exception as e:
        logger.error(f"gTTS failed: {e}")
        raise HTTPException(status_code=500, detail=f"Text-to-speech conversion failed: {str(e)}")

    return {
        "message": "Voice overview generated successfully.",
        "audio_file": audio_filename,
        "download_url": f"/voice-overview/{audio_filename}",
        "summary_text": summary_text,
        "file_name": payload.file_name,
    }


@app.get("/voice-overview/{filename}", summary="Download a generated voice overview MP3")
async def download_voice_overview(filename: str):
    # Sanitize filename to prevent path traversal
    safe_filename = Path(filename).name
    audio_path = VOICE_DIR / safe_filename
    if not audio_path.exists():
        raise HTTPException(status_code=404, detail=f"Audio file '{filename}' not found.")
    return FileResponse(
        path=str(audio_path),
        media_type="audio/mpeg",
        filename=safe_filename,
    )


# ─────────────────────────────────────────────
# CONCEPT MAP
# ─────────────────────────────────────────────

@app.post("/concept-map", response_model=ConceptMapResponse, summary="Generate a concept map from documents")
async def generate_concept_map(
    payload: ConceptMapRequest,
    x_user_profile: str = Header(..., alias="X-User-Profile")
):
    user_id = x_user_profile.strip()
    user_dir = UPLOAD_DIR / user_id
    if payload.file_name:
        file_path = user_dir / payload.file_name
        if not file_path.exists():
            raise HTTPException(status_code=404, detail=f"File '{payload.file_name}' not found for your profile.")
        try:
            text = extract_text(file_path)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Text extraction failed: {str(e)}")
        text_sample = text[:10000]
        label = payload.file_name
    else:
        all_docs = vector_store.get(where={"user": user_id})
        if not all_docs or not all_docs.get("documents"):
            raise HTTPException(
                status_code=404,
                detail="No documents indexed yet for your profile. Please upload a file first.",
            )
        text_sample = "\n\n".join(all_docs["documents"][:25])[:10000]
        label = "all uploaded documents for your profile"

    output_format = (payload.output_format or "mermaid").lower()

    if output_format == "json":
        prompt = f"""You are a knowledge-graph expert.
Analyse the document content below and produce a hierarchical concept map as valid JSON.
The JSON MUST follow this exact schema with no extra text or markdown:
{{
  "root": "<central topic>",
  "branches": [
    {{
      "topic": "<branch name>",
      "subtopics": ["<item1>", "<item2>", "<item3>"]
    }}
  ]
}}
Return ONLY the raw JSON — no markdown fences, no preamble, no extra text.

DOCUMENT CONTENT:
{text_sample}

JSON CONCEPT MAP:"""
    else:
        prompt = f"""You are a knowledge-graph expert.
Analyse the document content below and produce a valid Mermaid.js mindmap diagram.
Use EXACTLY this mindmap syntax with 2-space indentation:

mindmap
  root((Central Topic))
    Branch A
      Sub-item 1
      Sub-item 2
    Branch B
      Sub-item 3
      Sub-item 4

Rules:
- Start with "mindmap" on the first line
- The root node uses double parentheses: root((Topic Name))
- Top-level branches are indented 4 spaces
- Sub-items are indented 6 spaces
- Keep all node labels under 5 words and free of special characters like quotes or brackets
- Include 3-5 branches, each with 2-3 sub-items
- Return ONLY the raw mermaid text — no markdown fences, no extra text

DOCUMENT CONTENT:
{text_sample}

MERMAID MINDMAP:"""

    concept_map = call_llm(prompt)

    # Strip accidental markdown fences
    concept_map = concept_map.strip()
    for fence in ("```mermaid", "```json", "```"):
        if concept_map.startswith(fence):
            concept_map = concept_map[len(fence):].strip()
    if concept_map.endswith("```"):
        concept_map = concept_map[:-3].strip()

    logger.info(f"Concept map generated for: {label} (format={output_format})")
    return ConceptMapResponse(
        concept_map=concept_map,
        output_format=output_format,
        file_name=payload.file_name,
    )


# ─────────────────────────────────────────────
# YOUTUBE VIDEO RECOMMENDATIONS
# ─────────────────────────────────────────────

def _iso8601_duration_to_seconds(duration: str) -> int:
    pattern = r'PT(?:(\d+)H)?(?:(\d+)M)?(?:(\d+)S)?'
    match = re.match(pattern, duration)
    if not match:
        return 0
    hours   = int(match.group(1) or 0)
    minutes = int(match.group(2) or 0)
    seconds = int(match.group(3) or 0)
    return hours * 3600 + minutes * 60 + seconds


def _seconds_to_label(total_seconds: int) -> str:
    h = total_seconds // 3600
    m = (total_seconds % 3600) // 60
    s = total_seconds % 60
    if h:
        return f"{h}h {m}m {s}s"
    if m:
        return f"{m}m {s}s"
    return f"{s}s"


def _detect_topic_from_docs(user_id: str, file_name: Optional[str] = None) -> str:
    user_dir = UPLOAD_DIR / user_id
    if file_name:
        file_path = user_dir / file_name
        if not file_path.exists():
            raise HTTPException(status_code=404, detail=f"File '{file_name}' not found for your profile.")
        try:
            text = extract_text(file_path)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Text extraction failed: {str(e)}")
        text_sample = text[:6000]
    else:
        all_docs = vector_store.get(where={"user": user_id})
        if not all_docs or not all_docs.get("documents"):
            raise HTTPException(
                status_code=404,
                detail="No documents indexed yet for your profile. Provide a 'topic' directly or upload a file.",
            )
        text_sample = "\n\n".join(all_docs["documents"][:20])[:6000]

    prompt = f"""Read the following document excerpt and identify the single most specific, \
searchable topic it covers. Reply with ONLY a short search query (3-7 words, no punctuation, \
no quotes) that would return the best YouTube tutorials or explanations about this topic.

DOCUMENT EXCERPT:
{text_sample}

SEARCH QUERY:"""

    topic = call_llm(prompt).strip().strip('"').strip("'")
    logger.info(f"Auto-detected YouTube topic: '{topic}'")
    return topic


@app.post(
    "/youtube-videos",
    response_model=YouTubeSearchResponse,
    summary="Find YouTube videos related to your document's topic",
)
async def youtube_videos(
    payload: YouTubeSearchRequest,
    x_user_profile: str = Header(..., alias="X-User-Profile")
):
    user_id = x_user_profile.strip()
    if not YOUTUBE_API_KEY:
        raise HTTPException(
            status_code=503,
            detail="YOUTUBE_API_KEY is not configured. Add YOUTUBE_API_KEY to your .env file. "
                   "Get a free key at: https://console.developers.google.com/",
        )

    if payload.topic and payload.topic.strip():
        topic = payload.topic.strip()
    else:
        topic = _detect_topic_from_docs(user_id, payload.file_name)

    sort_map = {
        "views":     "viewCount",
        "relevance": "relevance",
        "rating":    "rating",
        "date":      "date",
    }
    order = sort_map.get((payload.sort_by or "views").lower(), "viewCount")

    video_duration = None
    if payload.max_duration_minutes is not None:
        if payload.max_duration_minutes <= 4:
            video_duration = "short"
        elif payload.max_duration_minutes <= 20:
            video_duration = "medium"
        else:
            video_duration = "long"

    search_url = "https://www.googleapis.com/youtube/v3/search"
    max_results = min(max(payload.max_results or 5, 1), 10)
    search_params: dict = {
        "part":       "snippet",
        "q":          topic,
        "type":       "video",
        "order":      order,
        "maxResults": max_results * 3,
        "key":        YOUTUBE_API_KEY,
        "safeSearch": "moderate",
        "videoEmbeddable": "true",
    }
    if payload.language_code:
        search_params["relevanceLanguage"] = payload.language_code
    if video_duration:
        search_params["videoDuration"] = video_duration

    try:
        search_resp = requests.get(search_url, params=search_params, timeout=10)
        search_resp.raise_for_status()
    except requests.RequestException as e:
        logger.error(f"YouTube search API error: {e}")
        raise HTTPException(status_code=502, detail=f"YouTube API request failed: {str(e)}")

    items = search_resp.json().get("items", [])
    if not items:
        return YouTubeSearchResponse(topic=topic, videos=[], total_returned=0)

    video_ids = ",".join(
        item["id"]["videoId"] for item in items if "videoId" in item.get("id", {})
    )
    details_params = {
        "part": "statistics,contentDetails,snippet",
        "id":   video_ids,
        "key":  YOUTUBE_API_KEY,
    }

    try:
        details_resp = requests.get(
            "https://www.googleapis.com/youtube/v3/videos",
            params=details_params, timeout=10
        )
        details_resp.raise_for_status()
    except requests.RequestException as e:
        logger.error(f"YouTube details API error: {e}")
        raise HTTPException(status_code=502, detail=f"YouTube details API failed: {str(e)}")

    details_by_id = {v["id"]: v for v in details_resp.json().get("items", [])}
    max_duration_secs = (payload.max_duration_minutes * 60) if payload.max_duration_minutes else None

    videos: list[YouTubeVideoItem] = []
    for item in items:
        vid_id = item["id"].get("videoId")
        if not vid_id or vid_id not in details_by_id:
            continue

        detail  = details_by_id[vid_id]
        stats   = detail.get("statistics", {})
        content = detail.get("contentDetails", {})
        snippet = detail.get("snippet", {})

        duration_secs = _iso8601_duration_to_seconds(content.get("duration", "PT0S"))
        if max_duration_secs is not None and duration_secs > max_duration_secs:
            continue

        thumbnails = snippet.get("thumbnails", {})
        thumbnail = (
            thumbnails.get("high", {}).get("url") or
            thumbnails.get("medium", {}).get("url") or
            thumbnails.get("default", {}).get("url") or ""
        )
        description = snippet.get("description", "")

        videos.append(YouTubeVideoItem(
            title=snippet.get("title", "Untitled"),
            channel=snippet.get("channelTitle", "Unknown"),
            url=f"https://www.youtube.com/watch?v={vid_id}",
            thumbnail=thumbnail,
            duration_seconds=duration_secs,
            duration_label=_seconds_to_label(duration_secs),
            view_count=int(stats.get("viewCount", 0)),
            like_count=int(stats.get("likeCount", 0)),
            published_at=snippet.get("publishedAt", ""),
            description_snippet=description[:200] + "..." if len(description) > 200 else description,
        ))

        if len(videos) >= max_results:
            break

    logger.info(f"YouTube search: topic='{topic}', returned={len(videos)}")
    return YouTubeSearchResponse(topic=topic, videos=videos, total_returned=len(videos))


# ─────────────────────────────────────────────
# RUN
# ─────────────────────────────────────────────

if __name__ == "__main__":
    logger.info("━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")
    logger.info("  Mini NotebookLM — Backend Starting Up  ")
    logger.info(f"  LLM Provider : {LLM_PROVIDER.upper()}")
    logger.info(f"  Embedding    : {EMBEDDING_MODEL}")
    logger.info(f"  Upload Dir   : {UPLOAD_DIR.resolve()}")
    logger.info(f"  Vector DB    : {CHROMA_DIR.resolve()}")
    logger.info(f"  Server URL   : http://{HOST}:{PORT}")
    logger.info(f"  API Docs     : http://localhost:{PORT}/docs")
    logger.info("━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")
    uvicorn.run(
        "main_new:app",
        host=HOST,
        port=PORT,
        reload=True,
        log_level="info",
    )
